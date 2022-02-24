#!python3
from pptx import Presentation
import argparse
parser = argparse.ArgumentParser()
parser.add_argument('file', type=str, nargs='+', help='powerpoint file')
args = parser.parse_args()
prs = Presentation(args.file[0])
with open(args.file[0][:len(args.file[0]) - 5] + '.txt', 'a') as text_file:
    for slide in prs.slides:
        if len(slide.notes_slide.notes_text_frame.text) != 0:
            text_file.write(slide.shapes.title.text + '\n' +
                            slide.notes_slide.notes_text_frame.text + '\n\n\n')