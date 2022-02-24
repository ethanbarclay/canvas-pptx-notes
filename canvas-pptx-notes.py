#!python3
from canvasapi import Canvas
from io import BytesIO
from pptx import Presentation
import browser_cookie3
import requests
import os

API_URL = 'API_URL'
API_KEY = 'API_KEY'
COURSE_ID = 0

canvas = Canvas('https://' + API_URL, API_KEY)
modules = canvas.get_course(COURSE_ID).get_modules()

print('downloading...')

# get canvas cookies from chrome and add to request
cj = browser_cookie3.chrome(domain_name=API_URL)
cookie_string = ''.join([cookie.name + "=" + cookie.value + ";" for cookie in cj])
headers_dict = {"cookie": cookie_string}

# find pptx in modules
for module in modules:
    items = module.get_module_items()
    for item in items:
        # download pptx
        if 'pptx' in item.title:
            download_url = 'https://psu.instructure.com/files/' + \
                str(item.content_id) + '/download?download_frd=1'
            # follow redirect to cdn location
            download_url = requests.get(download_url,
                             cookies=cj, headers=headers_dict).history[1].url
            # download and read notes from presentation
            prs = Presentation(BytesIO(requests.get(download_url).content))
            with open(item.title[:len(item.title) - 5] + '.txt', 'a') as text_file:
                for slide in prs.slides:
                    if len(slide.notes_slide.notes_text_frame.text) != 0:
                        text_file.write(slide.shapes.title.text + '\n' +
                                        slide.notes_slide.notes_text_frame.text + '\n\n\n')
            # delete empty files
            if os.stat(item.title[:len(item.title) - 5] + '.txt').st_size == 0: 
                os.remove(item.title[:len(item.title) - 5] + '.txt')
print('finished!')