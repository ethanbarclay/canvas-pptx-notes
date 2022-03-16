#!python3
from pptx import Presentation
import argparse
import os

parser = argparse.ArgumentParser()
parser.add_argument('--file', type=str, nargs='+', help='powerpoint file')
args = parser.parse_args()

def process_presentation(presentation, *args):
    try:
        if presentation == args.file: prs = Presentation(args.file[0])
        else: prs = Presentation(presentation)
    except: prs = Presentation(presentation)
    if (len(args) > 0):
        if not os.path.exists(args[0]): os.makedirs(args[0])
        path = args[0] + "/" + prs.core_properties.title + '.txt'
    else: path = prs.core_properties.title + '.txt'
    with open(path, 'a') as text_file:
        for slide in prs.slides:
            if len(slide.notes_slide.notes_text_frame.text) != 0:
                text_file.write(slide.shapes.title.text + '\n' +
                                slide.notes_slide.notes_text_frame.text + '\n\n\n')
    # delete empty files
    if os.stat(path).st_size == 0: os.remove(path)
        
process_presentation(args.file)